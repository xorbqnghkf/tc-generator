import streamlit as st
import json
import io
import tempfile
import os
import requests
from pathlib import Path

# Excel 생성용
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# 문서 읽기용
import pdfplumber
from pptx import Presentation
import docx

# ── 페이지 설정 ────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="AI Test Case Generator",
    page_icon="🤖",
    layout="wide"
)

# ── CSS 스타일 ─────────────────────────────────────────────────────────────
st.markdown("""
<style>
    .main-title {
        font-size: 2.2rem;
        font-weight: 800;
        color: #1F4E79;
        margin-bottom: 0.2rem;
    }
    .sub-title {
        font-size: 1rem;
        color: #666;
        margin-bottom: 2rem;
    }
    .step-box {
        background: #F0F7FF;
        border-left: 4px solid #2E75B6;
        padding: 1rem 1.2rem;
        border-radius: 6px;
        margin-bottom: 1rem;
    }
    .result-box {
        background: #E8F5E9;
        border-left: 4px solid #4CAF50;
        padding: 1rem 1.2rem;
        border-radius: 6px;
    }
    .warn-box {
        background: #FFF8E1;
        border-left: 4px solid #FFC107;
        padding: 1rem 1.2rem;
        border-radius: 6px;
    }
    .stButton > button {
        background-color: #1F4E79;
        color: white;
        font-weight: 600;
        border-radius: 8px;
        padding: 0.6rem 2rem;
        border: none;
        font-size: 1rem;
    }
    .stButton > button:hover {
        background-color: #2E75B6;
    }
</style>
""", unsafe_allow_html=True)

# ── 헤더 ───────────────────────────────────────────────────────────────────
st.markdown('<div class="main-title">🤖 AI Test Case Generator</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-title">요구사항 문서를 업로드하면 AI가 자동으로 Test Case Excel을 생성합니다 (Local AI 기반)</div>', unsafe_allow_html=True)
st.divider()

# ── 사이드바 ───────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("### ⚙️ Ollama 설정")

    ollama_host = st.text_input(
        "Ollama 서버 주소",
        value="http://localhost:11434",
        help="기본값: http://localhost:11434"
    )

    # 사용 가능한 모델 목록 자동 조회
    available_models = []
    try:
        res = requests.get(f"{ollama_host}/api/tags", timeout=3)
        if res.status_code == 200:
            available_models = [m["name"] for m in res.json().get("models", [])]
    except:
        pass

    if available_models:
        selected_model = st.selectbox(
            "사용할 모델",
            options=available_models,
            help="Ollama에 설치된 모델 목록"
        )
        st.success(f"✅ Ollama 연결됨 ({len(available_models)}개 모델)")
    else:
        selected_model = st.text_input(
            "모델명 직접 입력",
            value="llama3.1:8b",
            help="예: llama3.1:8b, llama3.1:14b"
        )
        st.warning("⚠️ Ollama 연결 안 됨\nOllama가 실행 중인지 확인해주세요.")
        st.code("ollama serve", language="bash")

    st.markdown("---")
    st.markdown("### 📋 사용 방법")
    st.markdown("""
1. Ollama 실행 확인
2. 요구사항 문서 업로드
3. 프로젝트명 입력
4. **TC 생성** 버튼 클릭
5. Excel 다운로드
    """)
    st.markdown("---")
    st.markdown("### 📁 지원 파일 형식")
    st.markdown("- 📄 PDF (.pdf)\n- 📊 PowerPoint (.pptx)\n- 📝 Word (.docx)\n- 📃 텍스트 (.txt)")
    st.markdown("---")
    st.markdown("### 💡 팁")
    st.markdown("""
- 문서가 클수록 분석 시간이 걸려요
- 요구사항이 명확할수록 TC 품질↑
- 생성 후 검토·보완 권장
    """)

# ── 문서 텍스트 추출 함수 ──────────────────────────────────────────────────
def extract_text(uploaded_file):
    ext = Path(uploaded_file.name).suffix.lower()
    text = ""

    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(uploaded_file.getvalue())
        tmp_path = tmp.name

    try:
        if ext == ".pdf":
            with pdfplumber.open(tmp_path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"

        elif ext == ".pptx":
            prs = Presentation(tmp_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                text += line + "\n"

        elif ext == ".docx":
            doc = docx.Document(tmp_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

        elif ext == ".txt":
            text = uploaded_file.getvalue().decode("utf-8")

    finally:
        os.unlink(tmp_path)

    return text.strip()

# ── 텍스트 청크 분할 함수 (긴 문서 처리) ───────────────────────────────────
def split_chunks(text, chunk_size=4000, overlap=200):
    """긴 문서를 청크로 분할 (12,000자 제한 문제 해결)"""
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            # 문장 단위로 자르기
            cut = text.rfind('\n', start, end)
            if cut == -1:
                cut = end
            else:
                cut += 1
        else:
            cut = len(text)
        chunks.append(text[start:cut])
        start = cut - overlap
    return chunks

# ── Ollama API 호출 함수 ───────────────────────────────────────────────────
def call_ollama(host, model, prompt, system_prompt=""):
    """Ollama API 호출"""
    messages = []
    if system_prompt:
        messages.append({"role": "system", "content": system_prompt})
    messages.append({"role": "user", "content": prompt})

    response = requests.post(
        f"{host}/api/chat",
        json={
            "model": model,
            "messages": messages,
            "stream": False,
            "options": {
                "temperature": 0.1,   # 낮을수록 JSON 형식 안정적
                "num_predict": 8000,
            }
        },
        timeout=300  # 5분 타임아웃
    )
    response.raise_for_status()
    return response.json()["message"]["content"]

# ── Ollama 기반 분석 함수 ──────────────────────────────────────────────────
def analyze_with_ollama(host, model, doc_text, project_name):
    system_prompt = """당신은 전문 소프트웨어/하드웨어 검증 엔지니어입니다.
요구사항 문서를 분석하여 체계적인 Test Case를 생성합니다.

반드시 아래 JSON 형식으로만 응답하세요. 다른 텍스트는 절대 포함하지 마세요.
마크다운 코드블록(```)도 사용하지 마세요. 순수 JSON만 반환하세요:

{
  "project": "프로젝트명",
  "sheets": [
    {
      "name": "시트명 (예: 1.기능기본동작)",
      "title": "시트 전체 제목",
      "testcases": [
        {
          "type": "CAT",
          "category_name": "카테고리명 (예: [ 기능 분류명 ])"
        },
        {
          "type": "TC",
          "tc_id": "TC-001",
          "category": "분류",
          "req_type": "필수 또는 선택 또는 개발",
          "title": "테스트 항목명",
          "procedure": "테스트 절차 (구체적으로)",
          "expected": "기대 결과",
          "priority": "High 또는 Medium 또는 Low",
          "auto": "가능 또는 부분 또는 불가",
          "note": "비고 (없으면 빈 문자열)"
        }
      ]
    }
  ]
}

규칙:
- 요구사항 1개당 TC 1개 이상 생성
- 카테고리별로 CAT 항목으로 그룹화
- req_type은 문서에서 (필수)/(선택)/(개발) 표기 기준으로 분류
- auto 필드: WebUI에서 Playwright로 자동화 가능하면 "가능", 부분 가능하면 "부분", 수동만 가능하면 "불가"
- TC ID는 시트별로 독립적으로 부여
- 모든 TC는 검증 가능한 구체적 절차로 작성"""

    # 문서를 청크로 분할해서 처리
    chunks = split_chunks(doc_text, chunk_size=4000)
    all_sheets = []
    sheet_counter = 1

    for i, chunk in enumerate(chunks):
        user_prompt = f"""다음 요구사항 문서를 분석하여 Test Case를 생성해주세요.

프로젝트명: {project_name}
파트: {i+1}/{len(chunks)}

=== 요구사항 문서 ===
{chunk}
===================

위 문서의 모든 요구사항을 빠짐없이 TC로 변환해주세요.
관련 항목끼리 시트를 나누어 구성해주세요."""

        raw = call_ollama(host, model, user_prompt, system_prompt)

        # JSON 파싱
        raw = raw.strip()
        if "```json" in raw:
            raw = raw.split("```json")[1].split("```")[0].strip()
        elif "```" in raw:
            raw = raw.split("```")[1].split("```")[0].strip()

        # JSON 시작/끝 찾기
        start_idx = raw.find('{')
        end_idx = raw.rfind('}') + 1
        if start_idx != -1 and end_idx > start_idx:
            raw = raw[start_idx:end_idx]

        chunk_data = json.loads(raw)

        # 시트 번호 재정렬 (청크별로 중복 방지)
        for sheet in chunk_data.get("sheets", []):
            sheet["name"] = f"{sheet_counter}.{sheet['name'].split('.',1)[-1].strip() if '.' in sheet['name'] else sheet['name']}"
            sheet_counter += 1
            all_sheets.append(sheet)

    return {
        "project": project_name,
        "sheets": all_sheets
    }

# ── Excel 생성 함수 ────────────────────────────────────────────────────────
def create_excel(data):
    wb = Workbook()

    HDR_BG  = "1F4E79"
    CAT_BG  = "2E75B6"
    ODD_BG  = "F5F9FF"
    EVEN_BG = "FFFFFF"
    DEV_BG  = "FFF2CC"
    AUTO_BG = "E8F5E9"   # 자동화 가능 행 색상

    thin = Side(style="thin", color="BFBFBF")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    def fill(hex_): return PatternFill("solid", fgColor=hex_)
    def center(): return Alignment(horizontal="center", vertical="center", wrap_text=True)
    def left(): return Alignment(horizontal="left", vertical="center", wrap_text=True)

    # 자동화 컬럼 추가
    COLS = ["TC ID","분류","요구사항 유형","테스트 항목","테스트 절차","기대 결과","우선순위","자동화","비고","결과(Pass/Fail)"]
    COL_W = [9, 15, 10, 35, 42, 28, 10, 10, 18, 14]

    first = True
    for sheet_data in data.get("sheets", []):
        if first:
            ws = wb.active
            ws.title = sheet_data["name"][:31]
            first = False
        else:
            ws = wb.create_sheet(sheet_data["name"][:31])

        for i, w in enumerate(COL_W, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

        # 제목
        ws.merge_cells("A1:J1")
        ws["A1"] = sheet_data.get("title", sheet_data["name"])
        ws["A1"].font = Font(bold=True, name="Arial", size=12, color="FFFFFF")
        ws["A1"].fill = fill(HDR_BG)
        ws["A1"].alignment = center()
        ws.row_dimensions[1].height = 26

        # 헤더
        for c, name in enumerate(COLS, 1):
            cell = ws.cell(row=2, column=c, value=name)
            cell.font = Font(bold=True, name="Arial", size=9, color="FFFFFF")
            cell.fill = fill(CAT_BG)
            cell.alignment = center()
            cell.border = border
        ws.row_dimensions[2].height = 20

        row = 3
        for tc in sheet_data.get("testcases", []):
            if tc["type"] == "CAT":
                ws.merge_cells(f"A{row}:J{row}")
                c = ws.cell(row=row, column=1, value=tc.get("category_name", ""))
                c.font = Font(bold=True, name="Arial", size=9, color="FFFFFF")
                c.fill = fill(CAT_BG)
                c.alignment = left()
                c.border = border
                ws.row_dimensions[row].height = 16
            else:
                is_dev = tc.get("req_type", "") == "개발"
                is_auto = tc.get("auto", "") == "가능"
                if is_dev:
                    bg = DEV_BG
                elif is_auto:
                    bg = AUTO_BG
                else:
                    bg = ODD_BG if row % 2 == 1 else EVEN_BG

                values = [
                    tc.get("tc_id",""),
                    tc.get("category",""),
                    tc.get("req_type","필수"),
                    tc.get("title",""),
                    tc.get("procedure",""),
                    tc.get("expected",""),
                    tc.get("priority","Medium"),
                    tc.get("auto","불가"),
                    tc.get("note",""),
                    ""
                ]
                for c, val in enumerate(values, 1):
                    cell = ws.cell(row=row, column=c, value=val)
                    cell.font = Font(bold=(c==4), name="Arial", size=9)
                    cell.fill = fill(bg)
                    cell.alignment = center() if c in [1,3,7,8,10] else left()
                    cell.border = border
                ws.row_dimensions[row].height = 42

            row += 1

        ws.freeze_panes = "A3"

    # Summary 시트
    ws_sum = wb.create_sheet("0.Summary", 0)
    ws_sum.column_dimensions["A"].width = 30
    ws_sum.column_dimensions["B"].width = 12
    ws_sum.column_dimensions["C"].width = 12
    ws_sum.column_dimensions["D"].width = 12
    ws_sum.column_dimensions["E"].width = 12
    ws_sum.column_dimensions["F"].width = 12

    ws_sum.merge_cells("A1:F1")
    ws_sum["A1"] = f"{data.get('project','프로젝트')} - Test Case Summary"
    ws_sum["A1"].font = Font(bold=True, name="Arial", size=13, color="FFFFFF")
    ws_sum["A1"].fill = fill(HDR_BG)
    ws_sum["A1"].alignment = center()
    ws_sum.row_dimensions[1].height = 28

    for c, h in enumerate(["시트명","필수/선택 TC","개발 TC","자동화 가능","자동화 불가","합계"], 1):
        cell = ws_sum.cell(row=2, column=c, value=h)
        cell.font = Font(bold=True, name="Arial", size=10, color="FFFFFF")
        cell.fill = fill(CAT_BG)
        cell.alignment = center()
        cell.border = border

    total_req, total_dev, total_auto, total_manual, total_all = 0, 0, 0, 0, 0
    r = 3
    for sheet_data in data.get("sheets", []):
        tcs = [t for t in sheet_data.get("testcases", []) if t["type"] == "TC"]
        req_cnt  = len([t for t in tcs if t.get("req_type") != "개발"])
        dev_cnt  = len([t for t in tcs if t.get("req_type") == "개발"])
        auto_cnt = len([t for t in tcs if t.get("auto") == "가능"])
        man_cnt  = len([t for t in tcs if t.get("auto") == "불가"])
        all_cnt  = len(tcs)
        total_req += req_cnt; total_dev += dev_cnt
        total_auto += auto_cnt; total_manual += man_cnt; total_all += all_cnt

        for c, val in enumerate([sheet_data["name"], req_cnt, dev_cnt, auto_cnt, man_cnt, all_cnt], 1):
            cell = ws_sum.cell(row=r, column=c, value=val)
            cell.font = Font(name="Arial", size=10)
            cell.fill = fill("D6E4F0")
            cell.alignment = center() if c > 1 else left()
            cell.border = border
        ws_sum.row_dimensions[r].height = 20
        r += 1

    for c, val in enumerate(["합 계", total_req, total_dev, total_auto, total_manual, total_all], 1):
        cell = ws_sum.cell(row=r, column=c, value=val)
        cell.font = Font(bold=True, name="Arial", size=10, color="FFFFFF")
        cell.fill = fill(HDR_BG)
        cell.alignment = center() if c > 1 else left()
        cell.border = border
    ws_sum.row_dimensions[r].height = 22

    ws_sum.merge_cells(f"A{r+2}:F{r+2}")
    ws_sum[f"A{r+2}"] = "※ 노란색 행 = 개발 요구사항 / 초록색 행 = 자동화 가능 항목 (Playwright)"
    ws_sum[f"A{r+2}"].font = Font(name="Arial", size=9, color="C00000")

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output

# ── 메인 UI ────────────────────────────────────────────────────────────────
col1, col2 = st.columns([1, 1])

with col1:
    st.markdown("### 📁 문서 업로드")
    uploaded_file = st.file_uploader(
        "요구사항 문서를 여기에 드래그하거나 클릭하여 업로드하세요",
        type=["pdf", "pptx", "docx", "txt"],
        help="PDF, PowerPoint, Word, 텍스트 파일 지원"
    )

    project_name = st.text_input(
        "📌 프로젝트명",
        placeholder="예: olleh GiGA WiFi Premium 장비평가",
        help="생성될 Excel 파일의 제목으로 사용됩니다"
    )

with col2:
    st.markdown("### ℹ️ 안내")
    st.markdown("""
<div class="step-box">
<b>🔄 동작 순서</b><br><br>
1️⃣ Ollama 실행 확인 (사이드바)<br>
2️⃣ 왼쪽에 문서 업로드<br>
3️⃣ 프로젝트명 입력<br>
4️⃣ TC 생성 버튼 클릭<br>
5️⃣ Excel 다운로드!
</div>
""", unsafe_allow_html=True)

    st.markdown("""
<div class="warn-box">
<b>⚠️ 주의사항</b><br><br>
• 로컬 AI라 Claude보다 속도가 느릴 수 있어요<br>
• 문서가 길면 청크로 나눠서 처리해요<br>
• 이미지로만 된 내용은 추출이 제한됩니다<br>
• 초록색 행 = 자동화 가능 항목
</div>
""", unsafe_allow_html=True)

st.divider()

# ── TC 생성 버튼 ───────────────────────────────────────────────────────────
ollama_ready = bool(available_models) or bool(selected_model)

if uploaded_file and project_name and ollama_ready:
    if st.button("🚀 TC 생성 시작", use_container_width=True):
        with st.status("AI가 문서를 분석하고 있어요...", expanded=True) as status:

            # 1. 텍스트 추출
            st.write("📖 문서 텍스트 추출 중...")
            try:
                doc_text = extract_text(uploaded_file)
                char_count = len(doc_text)
                chunk_count = len(split_chunks(doc_text, chunk_size=4000))
                st.write(f"✅ 텍스트 추출 완료 ({char_count:,} 글자 / {chunk_count}개 파트로 분석)")
            except Exception as e:
                st.error(f"문서 읽기 실패: {e}")
                st.stop()

            if not doc_text:
                st.error("문서에서 텍스트를 추출할 수 없어요. 이미지 PDF일 수 있습니다.")
                st.stop()

            # 2. Ollama 분석
            st.write(f"🤖 로컬 AI ({selected_model})가 요구사항 분석 중... (시간이 걸릴 수 있어요)")
            try:
                result = analyze_with_ollama(ollama_host, selected_model, doc_text, project_name)
                total_tc = sum(
                    len([t for t in s.get("testcases", []) if t["type"] == "TC"])
                    for s in result.get("sheets", [])
                )
                auto_tc = sum(
                    len([t for t in s.get("testcases", []) if t["type"] == "TC" and t.get("auto") == "가능"])
                    for s in result.get("sheets", [])
                )
                st.write(f"✅ 분석 완료! TC {total_tc}개 생성 (자동화 가능: {auto_tc}개)")
            except Exception as e:
                st.error(f"AI 분석 실패: {e}")
                st.write("💡 Ollama가 실행 중인지 확인해주세요: `ollama serve`")
                st.stop()

            # 3. Excel 생성
            st.write("📊 Excel 파일 생성 중...")
            try:
                excel_file = create_excel(result)
                st.write("✅ Excel 생성 완료!")
            except Exception as e:
                st.error(f"Excel 생성 실패: {e}")
                st.stop()

            status.update(label="✅ 완료!", state="complete")

        # 결과 표시
        st.markdown("""
<div class="result-box">
<b>🎉 TC 생성 완료!</b><br>
아래 버튼을 클릭하여 Excel 파일을 다운로드하세요.
</div>
""", unsafe_allow_html=True)

        st.download_button(
            label="📥 Excel TC 파일 다운로드",
            data=excel_file,
            file_name=f"{project_name}_TestCase.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True
        )

        # 미리보기
        with st.expander("📋 생성된 TC 미리보기"):
            for sheet in result.get("sheets", []):
                tcs = [t for t in sheet.get("testcases", []) if t["type"] == "TC"]
                auto_cnt = len([t for t in tcs if t.get("auto") == "가능"])
                st.markdown(f"**{sheet['name']}** — {len(tcs)}개 TC (자동화 가능: {auto_cnt}개)")
                for tc in tcs[:3]:
                    auto_tag = "🟢" if tc.get("auto") == "가능" else "🔴"
                    st.markdown(f"- {auto_tag} `{tc.get('tc_id','')}` {tc.get('title','')}")
                if len(tcs) > 3:
                    st.markdown(f"  _...외 {len(tcs)-3}개_")

elif not uploaded_file:
    st.info("📁 요구사항 문서를 업로드해주세요.")
elif not project_name:
    st.info("📌 프로젝트명을 입력해주세요.")
elif not ollama_ready:
    st.warning("⚙️ 사이드바에서 Ollama 모델을 설정해주세요.")
